"""
任务A基线模型对比汇报 PPT — 按 scientific-slides 方法论重制

设计理念 (scientific-slides SKILL.md):
  - Visual-First: 每页至少一个强视觉元素，60-70% 视觉 / 30-40% 文字
  - 每页一个核心想法 (One idea per slide)
  - 极简文字: 3-4 bullets, 4-6 words each
  - 充足白空间 (40-50%)
  - 现代色盘: 血流/生物医学 → Teal + Coral
  - Story Arc: Hook → Context → Approach → Results(40-50%) → Implications → Closure
  - 高对比 (7:1)
  - 变化布局 (full-figure / two-column / big-number / table)
  - 15 页：几何特征与数据管线合并为一页；含 xlsx 同步的区域 mask 对照表
"""

import os

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = os.path.dirname(os.path.abspath(__file__))
PROJ = os.path.abspath(os.path.join(BASE, "..", ".."))
PLOTS = os.path.join(PROJ, "outputs", "field", "plots")
OUT_PPT = os.path.join(BASE, "任务A基线模型对比汇报.pptx")
XLS_GEOM = os.path.join(BASE, "显示几何特征.xlsx")
XLS_REGION = os.path.join(BASE, "复杂区域命名.xlsx")

# ═══ Color Palette: Medical/Hemodynamics ═══════════════════════════════
# Teal & Coral, high contrast, colorblind-safe
BG_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BG_LIGHT  = RGBColor(0xF7, 0xFA, 0xFC)
C_PRIMARY = RGBColor(0x0A, 0x93, 0x96)  # teal
C_DARK    = RGBColor(0x1B, 0x2A, 0x4A)  # dark navy
C_ACCENT  = RGBColor(0xEE, 0x6C, 0x4D)  # coral
C_TEXT    = RGBColor(0x2C, 0x2C, 0x2C)  # charcoal
C_SUBTEXT = RGBColor(0x6C, 0x75, 0x7D)  # muted gray
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_SUCCESS = RGBColor(0x27, 0xAE, 0x60)  # green
C_WARN    = RGBColor(0xE6, 0x7E, 0x22)  # orange
C_CARD_BG = RGBColor(0xF0, 0xF9, 0xF9)  # very light teal
C_STRIPE  = RGBColor(0x0A, 0x93, 0x96)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 15

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ═══ Helpers ════════════════════════════════════════════════════════════

def bg(slide, color=BG_WHITE):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def rect(slide, l, t, w, h, fill, line=None, radius=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background() if not line else (s.line.fill.solid(), setattr(s.line.fill.fore_color, 'rgb', line))
    return s

def txt(slide, l, t, w, h, text, sz=24, color=C_TEXT, bold=False, align=PP_ALIGN.LEFT, font="Arial"):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color; p.font.bold = bold
    p.font.name = font; p.alignment = align
    return tb

def multi(slide, l, t, w, h, lines, default_sz=20, default_clr=C_TEXT):
    """lines: list of str or (text, size, color, bold)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            s, sz, clr, b = item, default_sz, default_clr, False
        else:
            s = item[0]; sz = item[1] if len(item)>1 else default_sz
            clr = item[2] if len(item)>2 else default_clr; b = item[3] if len(item)>3 else False
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = s; p.font.size = Pt(sz); p.font.color.rgb = clr
        p.font.bold = b; p.font.name = "Arial"; p.space_after = Pt(4)
    return tb

def img(slide, path, l, t, width=None, height=None):
    if not os.path.exists(path):
        txt(slide, l, t, Inches(3), Inches(0.4), f"[Missing: {os.path.basename(path)}]", 10, C_ACCENT)
        return None
    kw = {}
    if width:  kw['width']  = width
    if height: kw['height'] = height
    return slide.shapes.add_picture(path, l, t, **kw)

def page_num(slide, n):
    txt(slide, Inches(12.4), Inches(7.1), Inches(0.8), Inches(0.3),
        f"{n}/{TOTAL}", 10, C_SUBTEXT, align=PP_ALIGN.RIGHT)

def title_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, SLIDE_W, Inches(0.06), C_PRIMARY)
    txt(slide, Inches(0.8), Inches(0.3), Inches(10), Inches(0.6), title, 36, C_DARK, True)
    if subtitle:
        txt(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(0.35), subtitle, 16, C_SUBTEXT)

def circle_num(slide, x, y, num, color=C_PRIMARY, size=Inches(0.5)):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid(); c.fill.fore_color.rgb = color; c.line.fill.background()
    c.text_frame.paragraphs[0].text = str(num)
    c.text_frame.paragraphs[0].font.size = Pt(18)
    c.text_frame.paragraphs[0].font.color.rgb = C_WHITE
    c.text_frame.paragraphs[0].font.bold = True
    c.text_frame.paragraphs[0].font.name = "Arial"
    c.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_table(
    slide,
    l,
    t,
    w,
    h,
    data,
    header_fill=C_DARK,
    hl_row=None,
    font_sz=13,
    body_align=PP_ALIGN.CENTER,
    header_align=PP_ALIGN.CENTER,
):
    rows, cols = len(data), len(data[0])
    ts = slide.shapes.add_table(rows, cols, l, t, w, h)
    tbl = ts.table
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = str(data[r][c])
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_sz)
                p.font.name = "Arial"
                p.alignment = header_align if r == 0 else body_align
                p.space_before = Pt(2)
                p.space_after = Pt(2)
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                for p in cell.text_frame.paragraphs:
                    p.font.color.rgb = C_WHITE
                    p.font.bold = True
            elif hl_row and r == hl_row:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xD1, 0xFA, 0xE5)
                for p in cell.text_frame.paragraphs:
                    p.font.bold = True
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF8, 0xF9, 0xFA)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_WHITE
    return ts


def load_geom_and_region_tables():
    """从汇报材料目录下两个 xlsx 读取与幻灯片一致的表格（缺文件时降级为内置占位）。"""
    geom_tbl = [["特征", "定义 / 处理口径"]]
    if os.path.isfile(XLS_GEOM):
        dg = pd.read_excel(XLS_GEOM, sheet_name=0, header=None)
        for i in range(1, len(dg)):
            a, b = dg.iloc[i, 0], dg.iloc[i, 1]
            if pd.isna(a) and pd.isna(b):
                continue
            ga = "" if pd.isna(a) else str(a).strip()
            gb = "" if pd.isna(b) else str(b).strip()
            if ga:
                geom_tbl.append([ga, gb])
    else:
        geom_tbl.append(["(未找到 xlsx)", XLS_GEOM])

    reg_tbl = None
    if os.path.isfile(XLS_REGION):
        dr = pd.read_excel(XLS_REGION, sheet_name=0, header=0)
        reg_tbl = [dr.columns.astype(str).tolist()] + dr.fillna("—").astype(str).values.tolist()
    else:
        reg_tbl = [["(未找到 xlsx)", XLS_REGION, "", "", ""]]

    return geom_tbl, reg_tbl


GEOM_TBL, REGION_TBL_FULL = load_geom_and_region_tables()

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title (Hook)
# Bold teal background, large title, minimal info
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
rect(sl, 0, 0, SLIDE_W, SLIDE_H, C_DARK)
rect(sl, 0, 0, Inches(0.12), SLIDE_H, C_PRIMARY)

txt(sl, Inches(1.5), Inches(1.6), Inches(10), Inches(1.2),
    "任务A：基线模型对比", 52, C_WHITE, True)
txt(sl, Inches(1.5), Inches(3.0), Inches(10), Inches(0.5),
    "显式几何特征 → 血流动力学场重建精度", 24, C_PRIMARY)

rect(sl, Inches(1.5), Inches(4.0), Inches(3.5), Inches(0.003), C_SUBTEXT)

multi(sl, Inches(1.5), Inches(4.3), Inches(6), Inches(1.6), [
    ("MLP → GraphSAGE → Transformer → +Geometry", 18, RGBColor(0xBD,0xC3,0xC7)),
    ("4 模型 · 3 seeds · 统一评估口径", 15, C_SUBTEXT),
    ("2026.03.26", 14, C_SUBTEXT),
])

# key result card
rect(sl, Inches(9.2), Inches(4.5), Inches(3.5), Inches(2.2), RGBColor(0x14,0x3D,0x59))
txt(sl, Inches(9.4), Inches(4.6), Inches(3), Inches(0.4), "核心发现", 14, C_PRIMARY, True)
txt(sl, Inches(9.4), Inches(5.1), Inches(3), Inches(0.6), "1.161", 48, C_WHITE, True)
txt(sl, Inches(9.4), Inches(5.75), Inches(3), Inches(0.3), "RMSE_|v| 最优", 14, C_PRIMARY)
txt(sl, Inches(9.4), Inches(6.1), Inches(3), Inches(0.3), "vs MLP −27%  vs GS −15%", 13, RGBColor(0xBD,0xC3,0xC7))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 2 — Context: Why this matters
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "研究背景", "为什么做这件事？"); page_num(sl, 2)

# Left column: 3 key bullets (big, sparse)
multi(sl, Inches(0.8), Inches(1.6), Inches(5.5), Inches(5), [
    ("血流动力学参数 (WSS/OSI) 是预测", 22, C_TEXT, True),
    ("动脉瘤破裂风险的核心指标", 22, C_TEXT, True),
    ("", 10, C_TEXT),
    ("CFD 模拟耗时数小时", 20, C_TEXT),
    ("→ GNN 代理模型实现秒级推断", 20, C_TEXT),
    ("", 10, C_TEXT),
    ("核心问题", 22, C_ACCENT, True),
    ("显式几何特征能否提升重建精度？", 22, C_ACCENT, True),
])

# Right: simple visual flow
models = [("MLP", "无图结构", C_SUBTEXT), ("GraphSAGE", "+图结构", C_PRIMARY),
           ("Transformer", "+注意力", C_WARN), ("Trans+Geom", "+几何特征", C_SUCCESS)]
for i, (name, desc, clr) in enumerate(models):
    y = Inches(1.8) + i * Inches(1.3)
    rect(sl, Inches(7.5), y, Inches(4.8), Inches(1.0), C_WHITE)
    rect(sl, Inches(7.5), y, Inches(0.08), Inches(1.0), clr)
    circle_num(sl, Inches(7.8), y + Inches(0.2), i+1, clr)
    txt(sl, Inches(8.5), y + Inches(0.1), Inches(3.5), Inches(0.4), name, 20, C_DARK, True)
    txt(sl, Inches(8.5), y + Inches(0.5), Inches(3.5), Inches(0.3), desc, 14, C_SUBTEXT)
    if i < 3:
        txt(sl, Inches(9.5), y + Inches(0.95), Inches(1), Inches(0.3), "↓", 18, C_SUBTEXT, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 3 — Experimental Setup (Approach)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "实验设计"); page_num(sl, 3)

cards = [
    ("6,966", "Graphs", "Train 4860 / Val 648 / Test 1458", C_PRIMARY),
    ("split_AG", "按患者分层", "无数据泄漏 · 可复现", C_PRIMARY),
    ("3 seeds", "Mean ± Std", "可靠性保障", C_PRIMARY),
    ("15,000", "Nodes/Graph", "GUO_XI_JIANG 81 snapshots", C_PRIMARY),
]
for i, (num, title, desc, clr) in enumerate(cards):
    x = Inches(0.6) + i * Inches(3.15)
    y = Inches(1.5)
    rect(sl, x, y, Inches(2.9), Inches(2.2), C_WHITE)
    rect(sl, x, y, Inches(2.9), Inches(0.06), clr)
    txt(sl, x + Inches(0.2), y + Inches(0.3), Inches(2.5), Inches(0.7), num, 36, C_DARK, True, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.2), y + Inches(1.0), Inches(2.5), Inches(0.3), title, 16, C_TEXT, True, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.2), y + Inches(1.4), Inches(2.5), Inches(0.5), desc, 12, C_SUBTEXT, False, PP_ALIGN.CENTER)

# feature table
feat_data = [
    ["模型", "输入特征", "图结构", "几何特征"],
    ["MLP", "coord + t + BC", "✗", "✗"],
    ["GraphSAGE", "coord + t + BC + is_wall", "✓", "✗"],
    ["Transformer", "coord + t + BC + is_wall", "✓", "✗"],
    ["Trans+Geom", "coord + t + BC + is_wall + geom", "✓", "✓"],
]
add_table(sl, Inches(0.6), Inches(4.2), Inches(12.1), Inches(2.5), feat_data,
          header_fill=C_PRIMARY, hl_row=4, font_sz=14)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 4 — 显式几何特征 + 混合采样/预处理（合并 · 上下布局 · 无口播/创新点）
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(
    sl,
    "显式几何特征与数据管线",
    "xlsx 同步 · hybrid 采样 · 中心线几何 → 节点特征",
)
page_num(sl, 4)

flow = [
    ("原始表面点", C_SUBTEXT),
    ("混合降采样\nFPS + Random", C_PRIMARY),
    ("中心线 / VMTK\n几何基底", C_WARN),
    ("显式特征\n+ 图构建", C_SUCCESS),
]
fx0 = Inches(0.4)
y_flow = Inches(1.12)
bh_flow = Inches(0.74)
for i, (label, clr) in enumerate(flow):
    x = fx0 + i * Inches(3.05)
    rect(sl, x, y_flow, Inches(2.72), bh_flow, C_WHITE)
    rect(sl, x, y_flow, Inches(2.72), Inches(0.06), clr)
    txt(
        sl,
        x + Inches(0.1),
        y_flow + Inches(0.1),
        Inches(2.52),
        Inches(0.58),
        label,
        12,
        C_DARK,
        True,
        PP_ALIGN.CENTER,
    )
    if i < len(flow) - 1:
        txt(
            sl,
            x + Inches(2.74),
            y_flow + Inches(0.18),
            Inches(0.32),
            Inches(0.45),
            "→",
            20,
            clr,
            True,
            PP_ALIGN.CENTER,
        )

# 上下布局：流程 → 显式几何特征表 → 混合采样/标准化表（通栏）
_TBL_W = Inches(12.35)
_TBL_X = Inches(0.42)
y_geom = y_flow + bh_flow + Inches(0.1)
h_geom = Inches(2.78)
add_table(
    sl,
    _TBL_X,
    y_geom,
    _TBL_W,
    h_geom,
    GEOM_TBL,
    header_fill=C_PRIMARY,
    font_sz=11,
    body_align=PP_ALIGN.LEFT,
)

pipe_tbl = [
    ["环节", "策略", "用意"],
    ["全局 hybrid", "~20% FPS + ~80% random", "骨架覆盖 + 统计多样性"],
    ["壁面点超额", "50% FPS + 50% random", "控显存 · 保留壁面几何"],
    ["特征标准化", "曲率 z-score；弧长/半径归一", "跨病例可比 · 可复现"],
]
y_pipe = y_geom + h_geom + Inches(0.08)
h_pipe = Inches(2.05)
add_table(
    sl,
    _TBL_X,
    y_pipe,
    _TBL_W,
    h_pipe,
    pipe_tbl,
    header_fill=C_DARK,
    font_sz=12,
    body_align=PP_ALIGN.LEFT,
)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 5 — Main Results Table (Results)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "主结果总表", "3 seeds mean±std · 测试集 1,458 graphs"); page_num(sl, 5)

main_data = [
    ["", "RMSE_u↓", "RMSE_v↓", "RMSE_w↓", "RMSE_|v|↓", "RMSE_p↓", "R²_p↑", "时延(ms)", "显存(MB)"],
    ["MLP",         "0.974", "0.976", "0.974", "1.600", "0.658", "0.920", "0.54", "127"],
    ["GraphSAGE",   "0.931", "0.917", "0.850", "1.361", "0.734", "0.901", "2.35", "530"],
    ["Transformer", "0.934", "0.917", "0.848", "1.364", "0.706", "0.908", "6.95", "2183"],
    ["Trans+Geom",  "0.898", "0.852", "0.696", "1.161", "0.654", "0.921", "6.88", "2182"],
]
add_table(sl, Inches(0.5), Inches(1.5), Inches(12.3), Inches(2.8), main_data,
          header_fill=C_DARK, hl_row=4, font_sz=14)

# Big-number cards below
deltas = [
    ("−27.4%", "vs MLP\n全部因素联合", C_PRIMARY),
    ("−14.9%", "vs GraphSAGE\n几何特征的增量", C_PRIMARY),
    ("−14.9%", "vs Transformer (无geom)\n同backbone · 仅加几何", C_SUCCESS),
]
for i, (num, desc, clr) in enumerate(deltas):
    x = Inches(0.5) + i * Inches(4.15)
    rect(sl, x, Inches(4.7), Inches(3.9), Inches(1.8), C_WHITE)
    rect(sl, x, Inches(4.7), Inches(3.9), Inches(0.06), clr)
    txt(sl, x + Inches(0.3), Inches(4.9), Inches(1.8), Inches(0.8), num, 42, clr, True)
    lines = desc.split("\n")
    multi(sl, x + Inches(2.2), Inches(5.0), Inches(1.5), Inches(1.0),
          [(l, 13, C_TEXT if i2==0 else C_SUBTEXT, i2==0) for i2, l in enumerate(lines)])

# Bottom insight bar
rect(sl, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.5), C_DARK)
txt(sl, Inches(0.8), Inches(6.72), Inches(12), Inches(0.45),
    "图结构有效；但拉开差距主要靠显式几何特征", 16, C_WHITE, True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 6 — Big Numbers: Wall vs Interior (Results deep-dive)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "壁面 vs 内部：几何特征的差异化收益"); page_num(sl, 6)

# Three big-number cards — this is the "one idea" for this slide
big_cards = [
    ("壁面区域", "0.038", "RMSE_|v|", "MLP=0.318  GS=0.095", "壁面精度 8× 优于 MLP", C_SUCCESS),
    ("内部区域", "2.067", "RMSE_|v|", "MLP=2.692  GS=2.317", "主要瓶颈 · 需后续优化", C_WARN),
    ("压力场", "0.921", "R²_p", "MLP=0.920  GS=0.901", "全模型 > 0.90 · 整体可靠", C_PRIMARY),
]
for i, (region, val, metric, comp, insight, clr) in enumerate(big_cards):
    x = Inches(0.5) + i * Inches(4.15)
    y = Inches(1.5)
    rect(sl, x, y, Inches(3.9), Inches(5.0), C_WHITE)
    rect(sl, x, y, Inches(3.9), Inches(0.08), clr)
    txt(sl, x + Inches(0.3), y + Inches(0.3), Inches(3.3), Inches(0.4), region, 20, clr, True, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.3), y + Inches(1.0), Inches(3.3), Inches(1.0), val, 64, C_DARK, True, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.3), y + Inches(2.2), Inches(3.3), Inches(0.3), metric, 14, C_SUBTEXT, False, PP_ALIGN.CENTER)
    rect(sl, x + Inches(0.5), y + Inches(2.8), Inches(2.9), Inches(0.003), RGBColor(0xE0,0xE0,0xE0))
    txt(sl, x + Inches(0.3), y + Inches(3.0), Inches(3.3), Inches(0.4), comp, 12, C_SUBTEXT, False, PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.3), y + Inches(3.7), Inches(3.3), Inches(0.8), insight, 15, C_TEXT, True, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 7 — Scatter Plots (Full-figure)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "预测 vs 真值：散点对比"); page_num(sl, 7)

img(sl, os.path.join(PLOTS, "fig_A3_multimodel_scatter_vel_mag.png"),
    Inches(0.3), Inches(1.3), width=Inches(6.3))
img(sl, os.path.join(PLOTS, "fig_A3_multimodel_scatter_p.png"),
    Inches(6.8), Inches(1.3), width=Inches(6.3))

txt(sl, Inches(0.3), Inches(6.8), Inches(6.3), Inches(0.3),
    "速度模长 |v|：Trans+Geom 最贴对角线", 14, C_SUBTEXT, False, PP_ALIGN.CENTER)
txt(sl, Inches(6.8), Inches(6.8), Inches(6.3), Inches(0.3),
    "压力 p：各模型均表现良好", 14, C_SUBTEXT, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 8 — Per-case Boxplot (Full-figure)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "病例级稳定性：Per-case 箱线图"); page_num(sl, 8)

img(sl, os.path.join(PLOTS, "fig_A4_multimodel_per_case_boxplot.png"),
    Inches(2.0), Inches(1.3), width=Inches(9.3))

rect(sl, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8), C_CARD_BG)
txt(sl, Inches(0.8), Inches(6.4), Inches(12), Inches(0.6),
    "Trans+Geom：中位数最低 · IQR 最窄 · 跨病例最稳定", 18, C_PRIMARY, True)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 9 — Regional Evaluation (Two-column, 图为主)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "分区域评估", "区域 mask 定义见下一页对照表 · 四模型统一口径"); page_num(sl, 9)

img(sl, os.path.join(PLOTS, "fig_A5_multimodel_regional_bar_rmse_vel_mag.png"),
    Inches(0.3), Inches(1.3), width=Inches(6.3))
img(sl, os.path.join(PLOTS, "fig_A5_multimodel_regional_bar_rmse_p.png"),
    Inches(6.8), Inches(1.3), width=Inches(6.3))

rect(sl, Inches(0.4), Inches(6.22), Inches(12.4), Inches(0.9), C_CARD_BG)
txt(
    sl,
    Inches(0.65),
    Inches(6.35),
    Inches(11.9),
    Inches(0.65),
    "复杂区域（高曲率 / 近壁 / 分叉等）收益可解释 · 非标签分布取巧",
    16,
    C_PRIMARY,
    True,
)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 10 — 复杂区域 mask 命名（与 复杂区域命名.xlsx 同步）
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "复杂区域：mask 命名与定义", "评估分区与 graph 特征同源 · 参数可调"); page_num(sl, 10)
add_table(
    sl,
    Inches(0.35),
    Inches(1.28),
    Inches(12.6),
    Inches(5.75),
    REGION_TBL_FULL,
    header_fill=C_SUCCESS,
    font_sz=10,
    body_align=PP_ALIGN.LEFT,
)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 11 — Case Panel (Full-figure)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "典型病例：GUO_XI_JIANG #1120", "四模型空间误差面板"); page_num(sl, 11)

img(sl, os.path.join(PLOTS, "fig_A2_case_panel_result_features_merged-1120.png"),
    Inches(0.2), Inches(1.2), width=Inches(12.9))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 12 — Scalar + Vector (Full-figure)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "标量场 & 矢量场", "速度模长 · 压力 · 面内速度方向"); page_num(sl, 12)

img(sl, os.path.join(PLOTS, "fig_A2_GUO1120_vel_p.png"),
    Inches(0.15), Inches(1.15), width=Inches(7.5))
img(sl, os.path.join(PLOTS, "fig_A2_GUO1120_vec.png"),
    Inches(7.7), Inches(1.15), width=Inches(5.5))

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 13 — Efficiency (Results: cost)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "效率与部署可行性", "精度—时延 Pareto · 成本对比"); page_num(sl, 13)

img(sl, os.path.join(PLOTS, "fig_A7_pareto_rmse_vel_mag_vs_latency_mean_std.png"),
    Inches(0.3), Inches(1.3), width=Inches(5.5))
img(sl, os.path.join(PLOTS, "fig_A7_efficiency_bars_mean_std.png"),
    Inches(6.0), Inches(1.3), width=Inches(3.6))
img(sl, os.path.join(PLOTS, "fig_A7_fullcase_peak_memory_per_seed.png"),
    Inches(9.7), Inches(1.3), width=Inches(3.4))

eff_data = [
    ["", "参数量", "时延(ms)", "显存(MB)", "RMSE_|v|"],
    ["MLP", "35.7K", "0.54", "127", "1.600"],
    ["GraphSAGE", "101K", "2.35", "530", "1.361"],
    ["Transformer", "250K", "6.95", "2,183", "1.364"],
    ["Trans+Geom", "250K", "6.88", "2,182", "1.161"],
]
add_table(sl, Inches(0.5), Inches(5.5), Inches(12.3), Inches(1.7), eff_data,
          header_fill=C_DARK, hl_row=4, font_sz=13)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusions (Implications)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "基线阶段结论"); page_num(sl, 14)

conclusions = [
    ("1", "图结构有效", "MLP→GS −14.9%\n壁面 0.318→0.095", C_PRIMARY),
    ("2", "backbone 非瓶颈", "Transformer ≈ GraphSAGE\n无几何时无差异", C_WARN),
    ("3", "几何是关键", "→+Geom −14.9%\nw 分量 R² +67.7%", C_SUCCESS),
    ("4", "区域收益真实", "高曲率/分叉区\n优势显著", C_ACCENT),
]
for i, (num, title, desc, clr) in enumerate(conclusions):
    x = Inches(0.4) + i * Inches(3.2)
    rect(sl, x, Inches(1.5), Inches(3.0), Inches(3.5), C_WHITE)
    rect(sl, x, Inches(1.5), Inches(3.0), Inches(0.08), clr)
    circle_num(sl, x + Inches(1.2), Inches(1.8), num, clr, Inches(0.55))
    txt(sl, x + Inches(0.2), Inches(2.5), Inches(2.6), Inches(0.4), title, 20, C_DARK, True, PP_ALIGN.CENTER)
    lines = desc.split("\n")
    multi(sl, x + Inches(0.2), Inches(3.1), Inches(2.6), Inches(1.5),
          [(l, 15, C_TEXT) for l in lines], default_sz=15)

# Bottleneck warning
rect(sl, Inches(0.4), Inches(5.4), Inches(12.5), Inches(0.8), RGBColor(0xFE,0xF3,0xE2))
rect(sl, Inches(0.4), Inches(5.4), Inches(0.08), Inches(0.8), C_WARN)
txt(sl, Inches(0.7), Inches(5.5), Inches(12), Inches(0.3), "当前瓶颈", 16, C_WARN, True)
txt(sl, Inches(0.7), Inches(5.85), Inches(12), Inches(0.3),
    "内部流场误差 > 2.0 · 壁面 RMSE 低 ≠ WSS 准 · 需近壁梯度优化", 14, C_TEXT)

# ═══════════════════════════════════════════════════════════════════════
# SLIDE 15 — Next Steps (Closure)
# ═══════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6]); bg(sl)
title_bar(sl, "下一步计划", "Line A (内部精度) + Line W (壁面导向) 并行"); page_num(sl, 15)

# Line A
rect(sl, Inches(0.4), Inches(1.3), Inches(6.1), Inches(0.45), C_PRIMARY)
txt(sl, Inches(0.6), Inches(1.32), Inches(5.5), Inches(0.4), "Line A · 内部精度优化", 18, C_WHITE, True)

lineA = [
    ["优先级", "实验", "内容"],
    ["P0", "Opt-01", "target_weights [2,2,2,0.5]"],
    ["P0", "Opt-02", "Pre-Norm LayerNorm"],
    ["P0", "Opt-03", "Opt-01 + Opt-02 组合"],
    ["P1", "Opt-04/05", "hidden=256, layers=4"],
    ["P2", "Opt-07", "区域加权 loss"],
]
add_table(sl, Inches(0.4), Inches(1.85), Inches(6.1), Inches(2.6), lineA,
          header_fill=C_PRIMARY, font_sz=13)

# Line W
rect(sl, Inches(6.8), Inches(1.3), Inches(6.1), Inches(0.45), C_SUCCESS)
txt(sl, Inches(7.0), Inches(1.32), Inches(5.5), Inches(0.4), "Line W · 壁面导向 (WSS/OSI)", 18, C_WHITE, True)

lineW = [
    ["实验", "内容", "目标"],
    ["W01", "近壁区域加权", "WSS 梯度"],
    ["W02", "壁面梯度监督", "WSS 精度"],
    ["W03", "直接 WSS 监督", "端到端最优"],
    ["W04", "两阶段训练", "壁面精调"],
]
add_table(sl, Inches(6.8), Inches(1.85), Inches(6.1), Inches(2.2), lineW,
          header_fill=C_SUCCESS, font_sz=13)

# Timeline
rect(sl, Inches(0.4), Inches(5.0), Inches(12.5), Inches(2.2), C_WHITE)
rect(sl, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.06), C_DARK)
txt(sl, Inches(0.7), Inches(5.15), Inches(12), Inches(0.35), "推荐节奏", 16, C_DARK, True)

weeks = [
    ("第1周", "P0 五组实验", C_PRIMARY),
    ("第2周", "容量扩展+区域加权", C_WARN),
    ("第3周", "最小必要消融", C_SUBTEXT),
    ("并行", "Line W 壁面优化", C_SUCCESS),
]
for i, (wk, desc, clr) in enumerate(weeks):
    x = Inches(0.6) + i * Inches(3.1)
    r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.65), Inches(2.8), Inches(1.3)) if False else None
    r = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(5.65), Inches(2.8), Inches(1.3))
    r.fill.solid(); r.fill.fore_color.rgb = clr; r.line.fill.background()
    tf = r.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = wk
    p.font.size = Pt(16); p.font.color.rgb = C_WHITE; p.font.bold = True
    p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = desc
    p2.font.size = Pt(12); p2.font.color.rgb = C_WHITE
    p2.font.name = "Arial"; p2.alignment = PP_ALIGN.CENTER

# ═══ Save ═══════════════════════════════════════════════════════════════
prs.save(OUT_PPT)
print(f"✅ PPT saved: {OUT_PPT}")
print(f"   Slides: {len(prs.slides)}")
